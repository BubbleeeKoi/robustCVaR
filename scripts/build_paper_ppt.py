"""Build a Chinese PPTX deck for the robust CVaR/NCVaR paper.

The script uses PyMuPDF-extracted Figure 1 assets and python-pptx native
shapes to create a concise, evidence-led academic presentation.
"""

from __future__ import annotations

from pathlib import Path
import os
import shutil
import zipfile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ppt_output"
FIG_DIR = OUT / "assets" / "figures"
PPTX = OUT / "final_presentation_cn.pptx"
QA = OUT / "qa_report.md"
MANIFEST = OUT / "asset_manifest.md"
OUTLINE = OUT / "ppt_outline_cn.md"
EXPLANATION = OUT / "paper_explanation_cn.md"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 249, 247)
INK = RGBColor(31, 38, 48)
MUTED = RGBColor(93, 103, 114)
ACCENT = RGBColor(29, 92, 135)
ACCENT2 = RGBColor(152, 64, 58)
PALE = RGBColor(231, 238, 236)
GREEN = RGBColor(54, 126, 92)
YELLOW = RGBColor(202, 152, 54)
FONT_CN = "SimSun"
FONT_EN = "Times New Roman"


def pt(size: float):
    return Pt(size)


def apply_font(run, cn_font=FONT_CN, en_font=FONT_EN):
    """Set Latin and East Asian fonts explicitly in the PPTX XML."""
    run.font.name = cn_font
    rpr = run._r.get_or_add_rPr()
    for tag, typeface in (("a:latin", en_font), ("a:ea", cn_font), ("a:cs", en_font)):
        el = rpr.find(qn(tag))
        if el is None:
            el = OxmlElement(tag)
            rpr.append(el)
        el.set("typeface", typeface)


def apply_paragraph_font(paragraph, cn_font=FONT_CN, en_font=FONT_EN):
    for run in paragraph.runs:
        apply_font(run, cn_font=cn_font, en_font=en_font)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    apply_paragraph_font(p)
    if align is not None:
        p.alignment = align
    return box


def add_title(slide, text, subtitle=None):
    add_text(slide, text, 0.62, 0.34, 9.0, 0.55, size=26, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.92, 8.4, 0.34, size=9, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.18), Inches(1.25), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def add_source(slide, text):
    add_text(slide, text, 0.62, 7.08, 12.0, 0.22, size=7.2, color=MUTED)


def add_bullets(slide, bullets, x, y, w, h, size=15, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = pt(size)
        p.font.color.rgb = color
        apply_paragraph_font(p)
        p.level = 0
        p.space_after = pt(5)
    return box


def add_panel(slide, x, y, w, h, label, body, fill=PALE, color=INK, body_size=14):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(210, 217, 216)
    add_text(slide, label, x + 0.16, y + 0.12, w - 0.3, 0.25, size=10, color=ACCENT, bold=True)
    add_text(slide, body, x + 0.16, y + 0.45, w - 0.3, h - 0.55, size=body_size, color=color)
    return shp


def add_condition_card(slide, x, y, w, h, label, body):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PALE
    shp.line.color.rgb = RGBColor(210, 217, 216)

    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.15), Inches(w - 0.35), Inches(h - 0.25))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    p = tf.paragraphs[0]
    p.text = label
    p.font.size = pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    apply_paragraph_font(p)
    p.space_after = pt(5)

    p = tf.add_paragraph()
    p.text = body
    p.font.size = pt(12.4)
    p.font.color.rgb = INK
    apply_paragraph_font(p)
    p.space_after = pt(0)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.8)
    return conn


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def add_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = text


def make_fig1_contact_sheet():
    order = [
        ("a  no uncertainty", FIG_DIR / "fig1_panel_3.png"),
        ("b  RN: K=2", FIG_DIR / "fig1_panel_1.png"),
        ("c  KL: K=2", FIG_DIR / "fig1_panel_2.png"),
        ("d  decision-dependent: K in [1,2]", FIG_DIR / "fig1_panel_0.png"),
    ]
    thumbs = []
    for label, path in order:
        im = Image.open(path).convert("RGB")
        im = im.resize((545, 396))
        canvas = Image.new("RGB", (545, 430), "white")
        draw = ImageDraw.Draw(canvas)
        canvas.paste(im, (0, 34))
        draw.rectangle((0, 0, 545, 34), fill=(245, 247, 246))
        draw.text((10, 9), label, fill=(35, 45, 55))
        thumbs.append(canvas)
    sheet = Image.new("RGB", (1120, 900), "white")
    sheet.paste(thumbs[0], (0, 0))
    sheet.paste(thumbs[1], (575, 0))
    sheet.paste(thumbs[2], (0, 470))
    sheet.paste(thumbs[3], (575, 470))
    out = FIG_DIR / "fig1_contact_sheet.png"
    sheet.save(out, quality=95)
    return out


def build_deck():
    OUT.mkdir(exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    contact = make_fig1_contact_sheet()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slides = []
    for _ in range(16):
        slide = prs.slides.add_slide(blank)
        set_bg(slide)
        slides.append(slide)

    # 1
    s = slides[0]
    add_text(s, "鲁棒风险敏感强化学习：\n从 CVaR 到 NCVaR", 0.75, 0.8, 8.6, 1.2, 31, bold=True)
    add_text(s, "Robust Risk-Sensitive Reinforcement Learning with Conditional Value-at-Risk", 0.78, 2.08, 8.4, 0.35, 13, color=MUTED)
    add_panel(s, 0.78, 3.15, 3.1, 1.05, "这篇文章在做什么", "研究模型不确定下的\n风险敏感 RL", fill=PALE)
    add_panel(s, 4.15, 3.15, 3.1, 1.05, "核心工具", "固定预算：CVaR/EVaR\n决策相关：NCVaR", fill=PALE)
    add_panel(s, 7.52, 3.15, 3.1, 1.05, "验证场景", "64×53 GridWorld\n比较四种不确定性", fill=PALE)
    add_text(s, "Xinyi Ni, Lifeng Lai · IEEE ITW 2024 · DOI: 10.1109/ITW61385.2024.10806953", 0.78, 6.72, 11.5, 0.25, 8.5, color=MUTED)
    add_notes(s, "开场强调：本文不是普通 CVaR RL，而是问 CVaR 在模型不确定时是否仍然稳健，并进一步处理决策会改变不确定性的情况。")

    # 2
    s = slides[1]
    add_title(s, "一页讲清楚：这篇文章到底在做什么", "p.1 Introduction；p.2-p.5 Method and Experiment")
    add_panel(s, 0.85, 1.45, 5.15, 1.35, "研究对象", "在 MDP/RL 中，转移概率 P 往往来自估计。\n如果 P 错了，原来最优的策略可能在真实环境里很危险。", fill=PALE, body_size=12.2)
    add_panel(s, 6.75, 1.45, 5.15, 1.35, "风险目标", "文章不是只最小化平均代价，而是最小化 CVaR：\n关注最坏一部分轨迹中的累计代价。", fill=RGBColor(236, 242, 245), body_size=12.2)
    add_panel(s, 0.85, 3.25, 5.15, 1.42, "核心问题", "把“模型最坏情形”与“尾部风险最坏情形”放在一起：\nminπ maxP~ CVaRα[discounted cost]。", fill=RGBColor(244, 239, 228), body_size=12.2)
    add_panel(s, 6.75, 3.25, 5.15, 1.42, "解决路线", "固定不确定预算：转成 CVaR/EVaR。\n决策相关预算：提出 NCVaR，并给出 value iteration。", fill=RGBColor(242, 232, 228), body_size=12.2)
    add_text(s, "一句话：这篇文章把“模型不确定性”与“尾部风险控制”统一到一个可计算的风险敏感强化学习框架中。", 0.95, 5.48, 11.2, 0.42, 15.2, bold=True, color=ACCENT)
    add_source(s, "来源：p.1 Introduction；p.2 Eq.1-Eq.2；p.3 Section III；p.4 Theorem 1；p.5 Experiment")
    add_notes(s, "这一页是整篇文章的总括。先说对象：强化学习里转移概率不一定准；再说目标：不是平均最优，而是控制尾部损失；然后说难点：模型最坏情形和 CVaR 尾部最坏情形叠加；最后说解决路线：固定预算可化成 CVaR/EVaR，决策相关预算用 NCVaR 和动态规划处理。")

    # 3
    s = slides[2]
    add_title(s, "先统一语言：本文反复出现的 6 个术语", "p.1-p.4 definitions and setup")
    add_panel(s, 0.72, 1.35, 3.7, 1.05, "MDP / RL", "通俗：人在环境里连续做决策。\n专业：状态 x、动作 a、转移 P、代价 C。", fill=PALE, body_size=10.5)
    add_panel(s, 4.82, 1.35, 3.7, 1.05, "CVaRα", "通俗：只看最坏一部分结果。\n专业：分布尾部 α 分位的条件期望。", fill=RGBColor(236, 242, 245), body_size=10.5)
    add_panel(s, 8.92, 1.35, 3.7, 1.05, "RMDP", "通俗：承认模型会错。\n专业：在 uncertainty set 内做 worst-case 优化。", fill=PALE, body_size=10.5)
    add_panel(s, 0.72, 3.05, 3.7, 1.05, "ambiguity set", "通俗：真实模型允许偏离名义模型的范围。\n专业：由 divergence ball 定义的 P~ 集合。", fill=RGBColor(244, 239, 228), body_size=10.2)
    add_panel(s, 4.82, 3.05, 3.7, 1.05, "EVaR", "通俗：比 CVaR 更紧的尾部风险上界。\n专业：KL ball 对偶对应的风险度量。", fill=RGBColor(242, 232, 228), body_size=10.2)
    add_panel(s, 8.92, 3.05, 3.7, 1.05, "NCVaR / κ", "通俗：不同决策有不同不确定性。\n专业：κ(x,a) 进入风险度量并随路径变化。", fill=RGBColor(236, 242, 245), body_size=10.2)
    add_text(s, "听众只需抓住主线：模型可能错，坏结果很重要；本文用风险度量把两者合起来。", 0.9, 5.45, 11.5, 0.38, size=14.5, bold=True, color=ACCENT)
    add_source(s, "来源：p.1 Introduction；p.2 CVaR/EVaR；p.2 Eq.1 ambiguity set；p.3-p.4 NCVaR")
    add_notes(s, "这一页是面向所有人的术语翻译。对非专业听众，强调 MDP 就是连续决策，CVaR 就是只看坏结果平均有多坏，ambiguity set 就是真实模型允许偏离名义模型的范围。对专业听众，可以补充：RMDP 是对转移核做 worst-case，CVaR/EVaR 都可由对偶集合表征，κ(x,a) 是 decision-dependent uncertainty 的关键。")

    # 4
    s = slides[3]
    add_title(s, "为什么要做：已有两条路线各缺一半", "p.1 Introduction")
    add_panel(s, 0.8, 1.45, 3.45, 1.55, "路线 A：RMDP", "能处理模型不准，\n但多优化平均代价；\n坏结果有多坏不突出。", fill=PALE, body_size=11.6)
    add_panel(s, 4.9, 1.45, 3.45, 1.55, "路线 B：CVaR RL", "能处理尾部风险，\n但常假设转移模型 P 固定；\n模型误差不突出。", fill=PALE, body_size=11.6)
    add_panel(s, 9.0, 1.45, 3.45, 1.55, "本文缺口", "如果模型也不准、\n坏结果也重要，\n如何统一优化？", fill=RGBColor(242, 232, 228), body_size=11.6)
    add_bullets(s, ["面向非专业听众：这是“地图不准 + 事故代价高”时如何规划路线的问题。", "面向专业听众：这是 minπ maxP~ CVaRα 目标的等价表示与 Bellman 可解性问题。", "本文贡献不是只做 GridWorld，而是把固定预算和决策相关预算分别纳入风险度量。"], 0.95, 3.72, 11.2, 1.8, size=12.7)
    add_source(s, "来源：p.1 Introduction；p.3 Section III；p.4 Section IV")
    add_notes(s, "这里承接术语页。可以用一个直观例子讲：如果自动驾驶或机器人知道地图不完全准确，同时撞上障碍代价很高，那么平均最短路径不一定可靠。RMDP 只处理地图不准，CVaR 只处理坏结果，本文要把两者连起来。专业听众关心的是这个目标是否能化成可计算的动态规划，后面几页回答这个问题。")

    # 5
    s = slides[4]
    add_title(s, "问题形式：在所有可能转移中最小化最坏 CVaR", "p.2 Eq.1；p.2 Eq.2")
    add_text(s, "目标：minπ maxP~ CVaRα[Σ γᵗ C(xₜ,aₜ)]", 0.9, 1.48, 7.1, 0.55, 21, bold=True, color=ACCENT)
    add_bullets(s, ["P 是名义转移；P~ 是 ambiguity set 中可能真实发生的转移。", "φ-divergence 负责规定 P~ 离 P 可以有多远。", "CVaRα 只看累计代价分布中最坏 α 比例的尾部。"], 0.95, 2.25, 5.65, 1.75, size=12.8)
    add_panel(s, 7.15, 1.48, 4.75, 2.25, "逐层读公式", "minπ：我们选策略\nmaxP~：自然选择最坏模型\nCVaRα：只评价尾部坏结果\nΣγᵗC：折扣累计代价", fill=RGBColor(235, 242, 246), body_size=12.0)
    add_panel(s, 1.0, 4.75, 10.8, 0.86, "一句话", "论文优化的不是“平均情况下走得短”，而是“模型估错、坏结果发生时，策略仍尽量不付出极端代价”。", fill=RGBColor(244, 239, 228), body_size=12.6)
    add_source(s, "来源：p.2 Eq.1 ambiguity set；p.2 Eq.2 robust CVaR objective")
    add_notes(s, "这一页用于解释数学目标。强调 max 和 CVaR 的双重保守性。")

    # 6
    s = slides[5]
    add_title(s, "固定 RN 预算：鲁棒 CVaR 退化为更小置信水平的 CVaR", "p.3 Section III-A")
    add_text(s, "DRN(P~,P) ∈ [0,K]  ⇒  α′CVaR = α / K", 1.05, 1.48, 7.6, 0.52, 22, bold=True, color=ACCENT)
    add_bullets(s, ["RN derivative 控制概率密度比上界", "固定 K 可合并两层 worst-case", "Figure 1b 取 α=0.48, K=2，因此 α′=0.24"], 1.05, 2.35, 5.9, 1.8)
    add_panel(s, 7.55, 2.0, 3.7, 1.45, "直觉", "K 越大，等价 α′ 越小\n策略越关注更极端尾部", fill=RGBColor(236, 242, 245))
    add_source(s, "来源：p.3 Eq.5；p.5 Experiment")
    add_notes(s, "解释 Figure 1b 的设计：不是重新发明一个算法，而是利用等价关系读更小 α 的 CVaR 策略。")

    # 7
    s = slides[6]
    add_title(s, "固定 KL 预算：通过 EVaR 对偶处理指数尾部风险", "p.2 EVaR 定义；p.3 Section III-B")
    add_text(s, "DKL(Q,P) ≤ −ln α′  ⇔  EVaRα′", 0.95, 1.48, 7.2, 0.52, 22, bold=True, color=ACCENT)
    add_bullets(s, ["KL ambiguity 约束分布偏离程度", "EVaR 的对偶集合正是 KL ball", "Figure 1c 报告 α′EVaR=0.03"], 1.05, 2.38, 5.7, 1.65)
    add_panel(s, 7.35, 1.65, 4.1, 2.3, "为什么更保守", "KL 允许概率质量\n向高代价结果倾斜；\nEVaR 对尾部更敏感，\n通常比 CVaR 更保守。", fill=RGBColor(242, 232, 228), body_size=12.0)
    add_source(s, "来源：p.2 EVaR dual representation；p.3 Eq.6；p.5 α′EVaR=0.03")
    add_notes(s, "说明 1c 为什么和 1b 不同：RN 是密度比上界，KL 是相对熵球，风险度量也从 CVaR 变成 EVaR。")

    # 8
    s = slides[7]
    add_title(s, "为什么需要 NCVaR：不确定性会被决策本身改变", "p.3-p.4 Section IV")
    add_panel(s, 0.9, 1.45, 3.2, 1.45, "固定预算 K", "所有状态动作共用一个 K。\n最坏情形只会把 α 缩小，\n所以可化成普通 CVaR。", fill=PALE, body_size=11.5)
    add_panel(s, 4.9, 1.45, 3.2, 1.45, "决策相关 κ(x,a)", "某些动作更不确定，\n某些区域模型更不可靠。\n预算随策略路径变化。", fill=RGBColor(242, 232, 228), body_size=11.5)
    add_panel(s, 8.9, 1.45, 3.2, 1.45, "NCVaR 的作用", "把一串 κ(x,a) 放进\n风险度量本身，描述\n非均匀模型不确定性。", fill=PALE, body_size=11.5)
    add_bullets(s, ["论文假设 1 ≤ κ(x,a) ≤ Kmax：每个状态动作都有自己的不确定预算。", "这时不能简单写成 α/K，因为沿途每一步 K 都可能不同。", "Figure 1d 的 Kunfix∈[1,2] 就是让 κ 随状态动作变化，因此对应 NCVaR。"], 1.05, 3.65, 10.8, 1.7, size=12.8)
    add_source(s, "来源：p.3 Section IV；p.3 Assumption 1；p.5 Fig.1d")
    add_notes(s, "这里要突出 NCVaR 的必要性。固定 K 代表全环境同等不确定，因此可以直接调小 α。κ(x,a) 则表示不确定性由当前状态和动作决定，比如靠近障碍或采取激进动作时模型更不可靠。策略选择会影响未来遇到的 κ，所以必须把 κ 序列写进风险度量。")

    # 9
    s = slides[8]
    add_title(s, "如何解决：NCVaR 分解把轨迹风险变成 Bellman 更新", "p.4 Theorem 1；Bellman operator")
    add_text(s, "V(x,y) = minπ NCVaRy,κ [Σ γᵗ C | x₀=x]", 0.9, 1.45, 7.2, 0.45, 20, bold=True, color=ACCENT)
    add_bullets(s, ["把置信水平 y 纳入状态", "后继置信水平变成 y·ξ(x′)", "Bellman 中对 ξ 做 worst-case 最大化"], 1.05, 2.22, 5.6, 1.7)
    add_panel(s, 7.25, 1.45, 4.35, 2.25, "关键洞察", "风险预算不是一次性用完，而是在状态转移中递归分配；这使 NCVaR 可以做动态规划。", fill=RGBColor(236, 242, 245))
    add_panel(s, 1.05, 4.55, 10.6, 0.72, "讲解口径", "NCVaR 的贡献不是换一个名字，而是让 decision-dependent uncertainty 可递归计算。", fill=RGBColor(244, 239, 228))
    add_source(s, "来源：p.4 Theorem 1；p.4 Bellman operator T[V](x,y)")
    add_notes(s, "这一页是方法核心。强调 y 是增广维度，ξ 是密度比扰动，右侧置信水平随后继状态变。")

    # 10
    s = slides[9]
    add_title(s, "连续置信水平如何计算：对 yV(x,y) 做线性插值", "p.5 Eq.9-Eq.10；Algorithm 2")
    add_bullets(s, ["Y 是连续空间，直接遍历不可行", "采样 21 个 y 点并线性插值", "关键是插值 yV(x,y)，保持凹性和可计算性"], 0.95, 1.55, 4.65, 1.9)
    add_panel(s, 6.1, 1.35, 5.5, 2.9, "Algorithm 2", "1. 选择 Y(x), V0(x,y)\n2. 对所有 x,y 更新 TI[V]\n3. 收敛后用 V*(x,y) 构造 greedy policy", fill=RGBColor(236, 242, 245))
    add_panel(s, 0.95, 4.75, 10.65, 0.72, "为什么这么设计", "插值让 NCVaR 的连续置信维度变成有限维近似，同时保留 Bellman 算子的收敛结构。", fill=RGBColor(244, 239, 228))
    add_source(s, "来源：p.5 Eq.9 interpolation；p.5 Algorithm 2")
    add_notes(s, "解释实现细节：不是对 V 直接插值，而是对 yV 插值，这是论文为了保持凹性和 tractability 的关键。")

    # 11
    s = slides[10]
    add_title(s, "实验如何设计：用同一 GridWorld 比较四种不确定性", "p.5 Experiment")
    add_bullets(s, ["64×53 网格，80 个障碍物", "目标方向 0.95，其余方向 0.05/3", "安全移动代价 1，碰撞代价 40", "起点 (60,50)，终点 (60,2)"], 0.9, 1.4, 4.85, 2.5, size=14.2)
    add_condition_card(s, 6.05, 1.35, 2.95, 1.15, "1a", "CVaR\n无不确定性")
    add_condition_card(s, 9.35, 1.35, 2.95, 1.15, "1b", "RN K=2\nα′=0.24")
    add_condition_card(s, 6.05, 3.0, 2.95, 1.15, "1c", "KL K=2\nEVaR α′=0.03")
    add_condition_card(s, 9.35, 3.0, 2.95, 1.15, "1d", "K∈[1,2]\nNCVaR")
    add_source(s, "来源：p.5 Experiment；Fig.1 caption")
    add_notes(s, "回答实验为什么这样设计：同一个环境、同一 α，逐步改变不确定性集合，才能展示鲁棒性设定如何改变策略路径。")

    # 12
    s = slides[11]
    add_title(s, "核心证据：不确定性设定改变最优路径与风险分布", "p.5 Fig.1")
    s.shapes.add_picture(str(contact), Inches(0.78), Inches(1.25), width=Inches(7.35))
    add_bullets(s, ["蓝色表示低风险，黄色表示高风险", "红线为最优路径", "固定预算与决策相关预算均诱导更保守路径"], 8.55, 1.55, 3.75, 1.75, size=13.5)
    add_panel(s, 8.55, 4.2, 3.75, 1.1, "读图重点", "比较红线是否绕开\n障碍密集区；黄色\n风险区域是否扩大。", fill=RGBColor(244, 239, 228), body_size=12.2)
    add_source(s, "来源：p.5 Fig.1 原图；四个面板来自 PDF 第 5 页嵌入图像")
    add_notes(s, "提示听众按 a-b-c-d 顺序看：从无不确定性到固定预算，再到决策相关不确定性，路径整体更风险规避。")

    # 13
    s = slides[12]
    add_title(s, "验证逻辑：不是大规模 benchmark，而是机制性仿真实验", "p.5 Experiment")
    add_bullets(s, ["同一 GridWorld 控制环境变量", "同一 α=0.48 比较不同 ambiguity set", "路径和价值函数同时可视化", "目标是验证算法能诱导风险规避策略"], 0.9, 1.42, 5.0, 2.3)
    add_panel(s, 6.55, 1.45, 4.9, 1.25, "为什么这样设计", "GridWorld 简单但可解释：障碍物提供高代价尾部事件，路径变化能直接展示风险偏好。", fill=RGBColor(236, 242, 245))
    add_panel(s, 6.55, 3.25, 4.9, 1.25, "可信度边界", "论文没有公开障碍物坐标、γ、κ(x,a) 细节，因此 Figure 1 更像概念验证而非可完全复现实验。", fill=RGBColor(242, 232, 228))
    add_source(s, "来源：p.5 Experiment；Fig.1 caption；本文 PDF 未提供原始实验代码和 κ(x,a)")
    add_notes(s, "这页要讲清楚验证方式和局限：实验设计服务于说明方法行为，而不是展示大规模泛化能力。")

    # 14
    s = slides[13]
    add_title(s, "意义与讨论：NCVaR 把决策相关不确定性带入风险敏感 RL", "p.5 Conclusion")
    add_panel(s, 0.9, 1.5, 3.25, 1.35, "解决了什么", "模型不确定 + 尾部风险\n统一到可计算框架", fill=PALE)
    add_panel(s, 4.95, 1.5, 3.25, 1.35, "为什么重要", "真实决策会改变不确定性\n固定 K 不够表达", fill=PALE)
    add_panel(s, 9.0, 1.5, 3.25, 1.35, "开放问题", "κ 如何学习？\n如何扩展到大规模 RL？", fill=RGBColor(242, 232, 228))
    add_bullets(s, ["理论意义：建立 robust CVaR 与风险度量之间的桥梁", "算法意义：给出 NCVaR value iteration 与插值近似", "应用意义：为安全、金融、机器人等尾部风险场景提供框架"], 1.05, 3.75, 10.6, 1.6)
    add_source(s, "来源：p.5-p.6 Conclusion and Future Direction")
    add_notes(s, "结尾强调贡献和边界：本文最有价值的是问题重写和 NCVaR 框架，而不是 GridWorld 本身。")

    # 15
    s = slides[14]
    add_title(s, "被追问时怎么回答：专业问题的边界在哪里", "p.3-p.6 Method, Experiment, Conclusion")
    add_panel(s, 0.75, 1.42, 3.65, 1.42, "Q1：这是不是过度保守？", "是的，robust + CVaR 是双重保守。\n文章目标不是平均收益最大，\n而是控制模型误差下的尾部损失。", fill=PALE, body_size=10.8)
    add_panel(s, 4.85, 1.42, 3.65, 1.42, "Q2：κ(x,a) 怎么来？", "论文假设 κ 已给定，\n用于表达决策相关不确定性；\n学习 κ 是未来方向。", fill=RGBColor(242, 232, 228), body_size=10.8)
    add_panel(s, 8.95, 1.42, 3.65, 1.42, "Q3：为什么 GridWorld 足够？", "它不是大规模 benchmark，\n而是机制性实验：\n路径变化能直接看出风险偏好。", fill=RGBColor(236, 242, 245), body_size=10.8)
    add_panel(s, 0.75, 3.5, 3.65, 1.42, "Q4：EVaR 和 CVaR 差别？", "CVaR 看尾部平均；\nEVaR 是更紧上界，\n对应 KL ambiguity 的对偶。", fill=RGBColor(236, 242, 245), body_size=10.8)
    add_panel(s, 4.85, 3.5, 3.65, 1.42, "Q5：算法能扩展吗？", "当前是 tabular value iteration。\n连续状态和深度 RL 需要近似，\n论文把它列为开放方向。", fill=PALE, body_size=10.8)
    add_panel(s, 8.95, 3.5, 3.65, 1.42, "Q6：复现为什么困难？", "PDF 未公开障碍坐标、γ、\nκ(x,a) 和原始代码。\n只能做机制复现。", fill=RGBColor(244, 239, 228), body_size=10.8)
    add_source(s, "来源：p.3-p.5 方法与实验；p.5-p.6 Conclusion；PDF 未给出原始代码和完整参数")
    add_notes(s, "这一页是答辩页。不要把它讲成缺陷清单，而是讲清边界：本文的理论贡献在于目标重写和 NCVaR 分解；实验是概念验证；κ 的来源、连续状态扩展和大规模深度 RL 是自然后续问题。面对专业人士追问时，要把已证明的部分和未解决的部分分开。")

    # 16
    s = slides[15]
    add_title(s, "最后带走三句话：所有听众都应理解的主线", "p.1-p.6 full paper synthesis")
    add_panel(s, 0.95, 1.45, 3.45, 1.45, "第一句话", "普通 RL 常默认模型准、看平均；\n但现实里模型会错，坏结果很贵。", fill=PALE, body_size=11.2)
    add_panel(s, 4.95, 1.45, 3.45, 1.45, "第二句话", "本文把模型不确定和尾部风险\n合成 robust CVaR 问题，\n再用风险度量对偶化简。", fill=RGBColor(236, 242, 245), body_size=11.2)
    add_panel(s, 8.95, 1.45, 3.45, 1.45, "第三句话", "当不确定性随决策变化时，\n提出 NCVaR，并给出\n可计算的 value iteration。", fill=RGBColor(242, 232, 228), body_size=11.2)
    add_bullets(s, ["非专业理解：这是在“不确定地图 + 高事故成本”下找稳健路线。", "专业理解：这是 robust risk-sensitive MDP 的风险度量等价与动态规划问题。", "科研价值：把固定 ambiguity set 与 decision-dependent uncertainty 放进同一个叙事框架。"], 1.05, 4.05, 10.9, 1.55, size=12.8)
    add_source(s, "来源：全文综合；p.1 Introduction；p.3-p.4 Method；p.5 Experiment；p.5-p.6 Conclusion")
    add_notes(s, "结尾时不要再堆公式。用三句话收束：第一，问题来自现实中的模型误差和尾部损失；第二，固定不确定预算可以通过 CVaR/EVaR 对偶处理；第三，决策相关不确定性需要 NCVaR。这样非专业听众能抓住直觉，专业听众也能看到理论贡献在哪里。")

    # Add slide numbers
    for i, slide in enumerate(slides, start=1):
        add_text(slide, f"{i:02d}", 12.28, 7.08, 0.35, 0.2, size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)

    prs.save(PPTX)
    return PPTX


def audit_deck(path: Path):
    prs = Presentation(path)
    issues = []
    media_count = len([n for n in zipfile.ZipFile(path).namelist() if n.startswith("ppt/media/")])
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                issues.append(f"slide {idx}: shape out of bounds")
    return len(prs.slides), media_count, issues


def write_reports(pptx: Path):
    slides, media, issues = audit_deck(pptx)
    soffice = detect_soffice()
    rendered = sorted((OUT / "rendered").glob("slide_*.png"))
    render_note = (
        f"已使用 aspose.slides 在 portfolio 环境导出 {len(rendered)} 张 PNG 预览；预览图带 Aspose evaluation watermark，仅用于 QA，不写入 PPTX。"
        if rendered
        else "当前未完成；可在 portfolio 环境安装 aspose.slides 并设置 DYLD_LIBRARY_PATH=/opt/homebrew/lib 后导出预览。"
    )
    OUTLINE.write_text(
        "\n".join(
            [
                "# 中文 PPT 汇报提纲",
                "",
                "## 论文定位",
                "",
                "- 类型：方法 / 算法论文",
                "- 汇报逻辑：problem-to-solution",
                "- 主线：模型不确定下的风险敏感 RL → 固定 ambiguity set 的 CVaR/EVaR 等价 → 决策相关 uncertainty 的 NCVaR → GridWorld 验证",
                "",
                "## 中心问题",
                "",
                "标准 MDP 假设转移概率固定，但真实 RL 中模型估计误差普遍存在；风险中性 RMDP 又忽略尾部高代价事件。本文要解决的是：如何在模型不确定时最小化轨迹代价的最坏情形 CVaR。",
                "",
                "## 核心贡献",
                "",
                "1. 固定预算 ambiguity set 下，把 robust CVaR 转换为已有风险敏感 RL 问题。",
                "2. RN ambiguity 对应更小置信水平的 CVaR；KL ambiguity 对应 EVaR。",
                "3. 决策相关 uncertainty 下定义 NCVaR，并给出分解定理与 value iteration。",
                "4. 用 GridWorld 显示不同 uncertainty set 会改变最优路径和价值函数。",
                "5. 为混合听众补充术语地图和专业问答边界。",
                "",
                "## 幻灯片结构",
                "",
                "1. 标题与一句话定位",
                "2. 一页讲清楚文章在做什么",
                "3. 关键术语地图",
                "4. RMDP 与 CVaR RL 的缺口",
                "5. Robust CVaR 问题形式",
                "6. RN 固定预算如何变成 CVaR",
                "7. KL 固定预算如何连接 EVaR",
                "8. 为什么 1d 是 NCVaR",
                "9. NCVaR Bellman 分解",
                "10. 线性插值 Algorithm 2",
                "11. GridWorld 实验设计",
                "12. Figure 1 核心证据",
                "13. 验证逻辑和边界",
                "14. 意义、局限和讨论",
                "15. 专业追问与回答边界",
                "16. 面向所有听众的三句话总结",
            ]
        ),
        encoding="utf-8",
    )
    MANIFEST.write_text(
        "\n".join(
            [
                "# 图像资产清单",
                "",
                "| 文件 | 来源 | 用途 | 质量说明 |",
                "| --- | --- | --- | --- |",
                "| fig1_panel_0.png | PDF p.5 Figure 1d | 组成 Figure 1 总览 | 保留原图面板信息 |",
                "| fig1_panel_1.png | PDF p.5 Figure 1b | 组成 Figure 1 总览 | 保留原图面板信息 |",
                "| fig1_panel_2.png | PDF p.5 Figure 1c | 组成 Figure 1 总览 | 保留原图面板信息 |",
                "| fig1_panel_3.png | PDF p.5 Figure 1a | 组成 Figure 1 总览 | 保留原图面板信息 |",
                "| fig1_contact_sheet.png | 整理自 PDF p.5 Figure 1 | Slide 12 核心证据 | 重新排版但未修改科学数据 |",
            ]
        ),
        encoding="utf-8",
    )
    EXPLANATION.write_text(
        "\n".join(
            [
                "# 论文中文讲解稿",
                "",
                "## 这篇文章在做什么",
                "",
                "这篇文章研究的是模型不确定条件下的风险敏感强化学习。普通强化学习通常假设转移概率 P 已知或估计足够准确，然后优化期望累计回报或期望累计代价。但现实中 P 往往来自有限数据估计，可能偏离真实环境；同时，在安全、金融、机器人等问题中，真正关心的不是平均表现，而是坏情形、尾部损失和罕见高代价事件。",
                "",
                "因此，文章把两个问题合在一起：第一，模型可能错；第二，即使模型固定，也不能只看平均值，而要看 CVaR 所刻画的尾部风险。它的目标可以理解为：在所有可能真实转移模型中，找到一个让最坏尾部累计代价尽量小的策略。",
                "",
                "## 解决了什么问题",
                "",
                "文章解决的是 robust MDP 与 CVaR reinforcement learning 之间的结合问题。已有 RMDP 文献通常处理模型不确定性，但多以期望代价为目标；已有 CVaR RL 文献能处理尾部风险，但常假设转移模型固定。本文要处理的是二者同时存在时的目标定义、等价转换和动态规划算法。",
                "",
                "## 关键术语怎么讲给不同听众",
                "",
                "- MDP/RL：对非专业听众说成“人在环境里连续做选择”；对专业听众指出状态、动作、转移概率和代价函数。",
                "- CVaR：对非专业听众说成“只看最坏一部分结果平均有多坏”；对专业听众说明它是尾部条件期望。",
                "- RMDP：对非专业听众说成“地图可能不准时仍要规划路线”；对专业听众说明它在 ambiguity set 内做 worst-case 优化。",
                "- ambiguity set：对非专业听众说成“真实模型允许偏离估计模型的范围”；对专业听众说明它由 φ-divergence ball 定义。",
                "- EVaR：对非专业听众说成“比 CVaR 更保守的尾部风险上界”；对专业听众说明它和 KL divergence ball 的对偶相关。",
                "- NCVaR 与 κ(x,a)：对非专业听众说成“不同动作和位置的不确定性不同”；对专业听众说明 κ 是 decision-dependent uncertainty budget，并进入风险度量和 Bellman 分解。",
                "",
                "## 为什么要解决这个问题",
                "",
                "如果只做风险中性优化，策略可能为了缩短平均路径而靠近障碍或高风险区域；一旦模型估计误差使坏事件概率上升，实际代价会很高。如果只做普通 CVaR，但忽略模型不确定性，得到的策略仍可能对转移概率估计非常敏感。本文的意义在于把模型误差和尾部风险同时纳入优化目标。",
                "",
                "## 如何解决",
                "",
                "文章先把模型不确定性写成 ambiguity set：真实转移 P~ 不一定等于名义转移 P，但二者距离受 φ-divergence 约束。然后研究 minπ maxP~ CVaRα 的目标。",
                "",
                "对于固定不确定预算，文章利用风险度量的对偶形式，把 worst-case 转移模型吸收到风险度量中：RN ambiguity 可以转成更小置信水平的 CVaR；KL ambiguity 可以转成 EVaR。这说明某些 robust CVaR 问题可以复用已有风险敏感 RL 思路。",
                "",
                "对于决策相关不确定性，预算不再是一个固定 K，而是 κ(x,a)。这表示不同状态动作具有不同模型不确定性。例如某些区域数据少、模型更不可靠，或某些动作更容易产生转移误差。此时不能再用固定 α/K 描述风险水平，因此文章提出 NCVaR，把沿途变化的 κ 纳入风险度量，并证明其可以做 Bellman 分解。",
                "",
                "## 实验如何设计",
                "",
                "实验使用 64×53 GridWorld，设置 80 个障碍物。动作有随机性：目标方向概率 0.95，其余方向共享 0.05。安全移动代价为 1，碰撞障碍代价为 40。这个环境的设计目的不是追求大规模 benchmark，而是让高代价尾部事件可视化：障碍物提供风险源，路径是否绕开障碍密集区可以直接反映策略风险偏好。",
                "",
                "Figure 1 比较四种设定：1a 是无模型不确定性的 CVaR；1b 是 RN ambiguity，K=2，相当于更小置信水平的 CVaR；1c 是 KL ambiguity，对应 EVaR；1d 是 decision-dependent uncertainty，κ 在 [1,2] 内变化，对应 NCVaR。四幅图共同说明，不确定性设定会改变价值函数风险分布和最优路径。",
                "",
                "## 这篇文章的意义",
                "",
                "理论上，它把 robust CVaR、CVaR/EVaR 对偶和 decision-dependent uncertainty 联系起来。方法上，它为 NCVaR 给出 Bellman 分解和插值 value iteration。应用上，它提供了一种处理模型误差与尾部风险叠加问题的框架，适合安全控制、金融决策和机器人规划等场景。",
                "",
                "## 需要注意的边界",
                "",
                "论文实验是机制验证，不是完整可复现 benchmark。PDF 没有公开障碍物坐标、折扣因子 γ、κ(x,a) 的具体构造和原始代码。因此复现 Figure 1 时，应区分严格复现和概念复现：前者需要作者原始参数，后者只能验证论文机制和趋势。",
                "",
                "## 现场可能被问到的问题",
                "",
                "1. 这是不是过度保守？回答：是双重保守，但目标就是控制模型误差下的尾部损失，不是追求平均最短路径。",
                "2. κ(x,a) 怎么获得？回答：本文把 κ 作为给定预算来分析，如何从数据中学习 κ 是未来方向。",
                "3. GridWorld 是否足够验证？回答：它是机制性实验，用来展示路径和风险分布如何随 uncertainty set 改变，不是大规模泛化 benchmark。",
                "4. EVaR 与 CVaR 为什么都出现？回答：不同 ambiguity set 的对偶风险度量不同；RN 对应 CVaR 调整，KL 对应 EVaR。",
                "5. 算法能否扩展到深度 RL？回答：论文当前给出 tabular value iteration 和插值近似，深度近似和连续状态扩展仍是开放问题。",
            ]
        ),
        encoding="utf-8",
    )
    QA.write_text(
        "\n".join(
            [
                "# QA 报告",
                "",
                f"- PPTX 生成状态：成功",
                f"- 文件：`{pptx.name}`",
                f"- 幻灯片数量：{slides}",
                f"- 嵌入媒体数量：{media}",
                "- 语言：中文，保留必要英文术语、公式和变量名",
                "- 字体：PPTX XML 已指定中文宋体 SimSun，英文/数字/公式 Times New Roman",
                "- 结构：按混合听众演讲逻辑组织：总览、术语、缺口、方法、证据、边界、问答",
                "- 讲解稿：已生成 `paper_explanation_cn.md`，包含术语解释、论文逻辑和现场问答口径",
                "- 图像：使用 PDF 第 5 页 Figure 1 面板，并生成总览图",
                "- 来源标注：每页底部包含小字号页码、公式或 Figure 来源",
                "- Speaker notes：每页已写入简短中文讲解备注",
                "- 自审修正：避免长段落，使用短标题、少量 bullet、图证据优先布局",
                "- 结构检查：已用 python-pptx 重新打开并检查 shape 是否越界",
                f"- 越界问题：{'无' if not issues else '; '.join(issues)}",
                f"- `soffice` 检测：{'找到 ' + soffice if soffice else '未在 PATH、/Applications、/opt/homebrew/bin、/usr/local/bin 中找到可执行 soffice'}",
                f"- 渲染预览：{render_note}",
                "- 依赖说明：PPTX 由 python-pptx 生成；PNG 预览可用 aspose.slides + mono-libgdiplus 完成。LibreOffice/soffice 不是当前必需路径。",
            ]
        ),
        encoding="utf-8",
    )


def detect_soffice() -> str | None:
    candidates = [
        os.environ.get("SOFFICE_PATH"),
        shutil.which("soffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


if __name__ == "__main__":
    pptx = build_deck()
    write_reports(pptx)
    print(pptx)
